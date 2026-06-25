import io
import json
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import unquote

import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from pptx import Presentation
from pypdf import PdfReader

load_dotenv()


class ResearchAssistant:
    def __init__(self, history_file: Optional[str] = None) -> None:
        self.session = requests.Session()
        self.history_file = Path(history_file or os.path.join(os.path.dirname(__file__), "research_history.json"))
        self._ensure_history_file()

    def _ensure_history_file(self) -> None:
        if not self.history_file.exists():
            self.history_file.write_text("[]", encoding="utf-8")

    def _load_history(self) -> List[Dict[str, Any]]:
        try:
            return json.loads(self.history_file.read_text(encoding="utf-8"))
        except Exception:
            return []

    def _save_history(self, history: List[Dict[str, Any]]) -> None:
        self.history_file.write_text(json.dumps(history, indent=2), encoding="utf-8")

    def save_run(self, topic: str, result: Dict[str, Any]) -> None:
        history = self._load_history()
        history.append(
            {
                "timestamp": datetime.utcnow().isoformat(),
                "topic": topic,
                "summary": result.get("summary", ""),
                "citations": result.get("citations", []),
                "workflow_steps": result.get("workflow_steps", []),
            }
        )
        self._save_history(history)

    def get_history(self) -> List[Dict[str, Any]]:
        return self._load_history()

    def search_web(self, query: str, max_results: int = 5) -> List[Dict[str, str]]:
        url = "https://html.duckduckgo.com/html/"
        try:
            response = self.session.get(
                url,
                params={"q": query},
                timeout=10,
                headers={"User-Agent": "Mozilla/5.0"},
            )
            response.raise_for_status()
        except requests.RequestException:
            return []

        soup = BeautifulSoup(response.text, "html.parser")
        results: List[Dict[str, str]] = []
        for result in soup.select(".result")[:max_results]:
            title_tag = result.select_one(".result__title a")
            snippet_tag = result.select_one(".result__snippet")
            if not title_tag:
                continue
            href = title_tag.get("href", "")
            if href.startswith("//"):
                href = "https:" + href
            results.append(
                {
                    "title": title_tag.get_text(" ", strip=True),
                    "url": unquote(href),
                    "snippet": snippet_tag.get_text(" ", strip=True) if snippet_tag else "",
                }
            )
        return results

    def read_pdf_text(self, pdf_bytes: Optional[bytes]) -> str:
        if not pdf_bytes:
            return ""
        try:
            reader = PdfReader(io.BytesIO(pdf_bytes))
            pages = [page.extract_text() or "" for page in reader.pages[:4]]
            return "\n".join(pages)
        except Exception:
            return ""

    def _call_openai(self, prompt: str, api_key: str) -> str:
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }
        payload = {
            "model": "gpt-4o-mini",
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.3,
        }
        response = requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers=headers,
            json=payload,
            timeout=30,
        )
        response.raise_for_status()
        data = response.json()
        return data["choices"][0]["message"]["content"]

    def _call_gemini(self, prompt: str, api_key: str) -> str:
        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {"temperature": 0.3},
        }
        response = requests.post(
            f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}",
            json=payload,
            timeout=30,
        )
        response.raise_for_status()
        data = response.json()
        return data["candidates"][0]["content"]["parts"][0]["text"]

    def summarize_text(self, text: str, topic: str, provider: str, api_key: str) -> str:
        if provider == "OpenAI" and api_key:
            try:
                prompt = (
                    f"Create a concise research summary for the topic '{topic}'. "
                    f"Use the following sources and notes:\n\n{text[:6000]}"
                )
                return self._call_openai(prompt, api_key)
            except Exception:
                pass
        if provider == "Gemini" and api_key:
            try:
                prompt = (
                    f"Create a concise research summary for the topic '{topic}'. "
                    f"Use the following sources and notes:\n\n{text[:6000]}"
                )
                return self._call_gemini(prompt, api_key)
            except Exception:
                pass

        if not text.strip():
            return (
                f"No source material was found for '{topic}'. "
                "Add a PDF or refine the topic to gather stronger context."
            )

        sentences = re.split(r"(?<=[.!?])\s+", text)
        top_sentences = [s.strip() for s in sentences if s.strip()][:6]
        summary = " ".join(top_sentences)
        return summary[:1800] if len(summary) > 1800 else summary

    def extract_key_points(self, text: str, topic: str, provider: str, api_key: str) -> List[str]:
        if provider == "OpenAI" and api_key:
            try:
                prompt = (
                    f"Extract 4 concise key points about '{topic}' from this research context. "
                    f"Return each point on a new line without bullets.\n\n{text[:4000]}"
                )
                response = self._call_openai(prompt, api_key)
                return [line.strip() for line in response.splitlines() if line.strip()][:4]
            except Exception:
                pass
        if provider == "Gemini" and api_key:
            try:
                prompt = (
                    f"Extract 4 concise key points about '{topic}' from this research context. "
                    f"Return each point on a new line without bullets.\n\n{text[:4000]}"
                )
                response = self._call_gemini(prompt, api_key)
                return [line.strip() for line in response.splitlines() if line.strip()][:4]
            except Exception:
                pass

        words = re.split(r"[^a-zA-Z0-9]+", text.lower())
        filtered = [w for w in words if len(w) > 4 and w not in {"topic", "about", "their", "there", "would", "could"}]
        unique_words = list(dict.fromkeys(filtered))[:6]
        return [f"The research suggests a strong focus on {word} in the context of {topic}." for word in unique_words]

    def clarify_query(self, query: str) -> str:
        text = (query or "").strip()
        if not text:
            return "I can help with that. Please tell me the topic you want to research."

        words = re.findall(r"[A-Za-z0-9]+", text.lower())
        if len(words) <= 3:
            return (
                "I can help with that, but I’d like a bit more detail. "
                "Which topic do you want to research, and do you want a summary, sources, citations, or a presentation?"
            )

        if any(term in text.lower() for term in ["tell me about", "about", "help me", "explain"]):
            if len(words) <= 5:
                return (
                    "I can help with that. Please share the exact topic and what you want from it, "
                    "such as a short summary, deeper explanation, or source list."
                )

        return ""

    def answer_followup(self, topic: str, question: str, context: str, provider: str, api_key: str) -> str:
        if provider == "OpenAI" and api_key:
            try:
                prompt = (
                    f"You are a research assistant. Answer the user's follow-up question based on the research context. "
                    f"Topic: {topic}\nQuestion: {question}\nContext: {context[:5000]}"
                )
                return self._call_openai(prompt, api_key)
            except Exception:
                pass
        if provider == "Gemini" and api_key:
            try:
                prompt = (
                    f"You are a research assistant. Answer the user's follow-up question based on the research context. "
                    f"Topic: {topic}\nQuestion: {question}\nContext: {context[:5000]}"
                )
                return self._call_gemini(prompt, api_key)
            except Exception:
                pass

        return (
            f"Based on the available research for '{topic}', I can give you a concise explanation or a deeper analysis. "
            f"Tell me what you want next, such as a summary, comparison, or source-backed answer."
        )

    def create_presentation(self, title: str, sections: List[Dict[str, str]]) -> bytes:
        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = "AI Research Assistant"

        for section in sections:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = section["title"]
            body = slide.placeholders[1].text_frame
            body.clear()
            body.text = section["content"]

        buffer = io.BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    def generate_citations(self, web_results: List[Dict[str, str]]) -> List[str]:
        citations = []
        for item in web_results:
            citations.append(f"{item['title']}. Available at: {item['url']}")
        return citations

    def _build_fallback_sources(self, topic: str) -> List[Dict[str, str]]:
        return [
            {
                "title": f"{topic} overview",
                "url": f"https://example.com/search?q={topic.replace(' ', '+')}",
                "snippet": f"A general reference overview for {topic} that can be used while live search results are unavailable.",
            }
        ]

    def run_workflow(
        self,
        topic: str,
        pdf_bytes: Optional[bytes] = None,
        provider: str = "Local fallback",
        api_key: str = "",
    ) -> Dict[str, Any]:
        web_results = self.search_web(topic)
        if not web_results:
            web_results = self._build_fallback_sources(topic)
        pdf_text = self.read_pdf_text(pdf_bytes)
        combined_text = "\n\n".join(
            [
                f"Topic: {topic}",
                *(f"Web result: {item['title']} - {item['snippet']}" for item in web_results[:3]),
                *( [f"PDF excerpt: {pdf_text[:2000]}"] if pdf_text else []),
            ]
        )
        summary = self.summarize_text(combined_text, topic, provider, api_key)
        key_points = self.extract_key_points(combined_text, topic, provider, api_key)
        sections = [
            {
                "title": "Research Overview",
                "content": summary[:1000],
            },
            {
                "title": "Key Sources",
                "content": "\n".join(f"- {item['title']}\n  {item['url']}" for item in web_results[:4]),
            },
            {
                "title": "Key Points",
                "content": "\n".join(f"- {point}" for point in key_points),
            },
        ]
        presentation_bytes = self.create_presentation(f"{topic} Research", sections)
        workflow_steps = [
            "Research collection",
            "Context extraction",
            "Summary generation",
            "Presentation preparation",
        ]
        result = {
            "summary": summary,
            "web_results": web_results,
            "citations": self.generate_citations(web_results),
            "presentation_bytes": presentation_bytes,
            "key_points": key_points,
            "context": combined_text,
            "workflow_steps": workflow_steps,
        }
        self.save_run(topic, result)
        return result
